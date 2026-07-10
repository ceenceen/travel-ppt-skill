# -*- coding: utf-8 -*-
"""travel-ppt 泛化生成器：把结构化行程数据 -> 16:9 可编辑 PPTX（每天单独一页，彩色背景图）。

用法:
  python gen_pptx.py <数据模块.py> <输出.pptx> [--photos DIR] [--budget XLSX] [--cover-slug NAME]

数据模块需定义: days(list, 必填) / hotels / timeline / budget / photo_map / TRIP_TITLE / TRIP_SUB / CLOSING
"""
import os, sys, io, re, glob, argparse
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------- 配色（已定稿：无橘色装饰块） ----------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0xB8, 0xC2, 0xCC)
BAR   = RGBColor(0x9A, 0xA8, 0xB4)   # 预算条形中性灰
WARN  = RGBColor(0xEA, 0xD9, 0xA0)   # 提醒淡金（非橘）
ST_OK = RGBColor(0x6F, 0xC2, 0x8B)   # 已订绿
ST_REC= RGBColor(0x7F, 0xC8, 0xD8)   # 推荐青

# 内置城市 -> 图片 slug 映射（可被数据模块 photo_map 覆盖）
CITY_PHOTO = {
    "哈尔滨":"harbin","长白山":"changbai","沈阳":"shenyang","丹东":"dandong",
    "北京":"beijing","天津":"tianjin","济南":"jinan","泰山":"taishan","泰安":"taishan",
    "深圳":"cover","上海":"shanghai","杭州":"hangzhou","西安":"xian","成都":"chengdu",
    "重庆":"chongqing","广州":"guangzhou","南京":"nanjing","武汉":"wuhan","青岛":"qingdao",
    "厦门":"xiamen","昆明":"kunming","大理":"dali","丽江":"lijiang","桂林":"guilin",
    "拉萨":"lhasa","敦煌":"dunhuang","张家界":"zhangjiajie","苏州":"suzhou","三亚":"sanya",
    "香港":"hongkong","返深":"cover",
}

def load_data(path):
    import importlib.util
    spec = importlib.util.spec_from_file_location("trip_data_mod", path)
    mod = importlib.util.module_from_spec(spec); spec.loader.exec_module(mod)
    return mod

# ---------------- helpers ----------------
def set_cjk(run, name="Microsoft YaHei"):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin","a:ea","a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)

def cover_img(slide, path):
    im = Image.open(path).convert("RGB"); w,h = im.size
    if w > 1920: im = im.resize((1920, int(h*1920/w)), Image.LANCZOS)
    buf = io.BytesIO(); im.save(buf, format="JPEG", quality=86); buf.seek(0)
    pic = slide.shapes.add_picture(buf, 0, 0)
    iw, ih = pic.width, pic.height
    scale = max(SW/iw, SH/ih)
    pic.width = int(iw*scale); pic.height = int(ih*scale)
    pic.left = int((SW-pic.width)/2); pic.top = int((SH-pic.height)/2)

def overlay(slide, opacity):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0,0,0); sh.line.fill.background()
    srgb = sh.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int((1-opacity)*65535))}))
    sh.shadow.inherit = False
    return sh

def panel(slide, left, top, width, height, opacity=0.52):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0,0,0); sh.line.fill.background()
    srgb = sh.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int((1-opacity)*65535))}))
    sh.shadow.inherit = False
    return sh

def text(slide, left, top, width, height, txt, size=28, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.12):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = txt
    run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    set_cjk(run); p.line_spacing = spacing
    return tb

def header_line(slide, left, top, txt, size=13, color=GRAY):
    return text(slide, left, top, Inches(6), Inches(0.4), txt.upper(), size=size, bold=True, color=color, spacing=1.0)

def col_bullets(slide, left, top, width, height, items, header, size=14.5, gap=7, hcolor=WHITE):
    text(slide, left, top, width, Inches(0.4), header, size=15, bold=True, color=hcolor)
    tb = slide.shapes.add_textbox(left, top+Inches(0.45), width, height-Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.12
        run = p.add_run(); run.text = "—  " + it
        run.font.size = Pt(size); run.font.color.rgb = WHITE; set_cjk(run)
    return tb

def new_slide(path, opacity=0.5):
    s = prs.slides.add_slide(BLANK)
    if path: cover_img(s, path)
    overlay(s, opacity)
    return s

def shorten(s, n=46):
    s = re.sub(r"\s+", " ", s).strip()
    return s if len(s) <= n else s[:n-1] + "…"

def photo(name):
    p = os.path.join(PHOTO, name + ".jpg")
    return p if os.path.exists(p) else (sorted(glob.glob(os.path.join(PHOTO,"*.jpg"))) or [None])[0]

def day_photo(d, photo_map):
    blob = " ".join(d.get("itinerary", [])) + " ".join(n for n,_ in d.get("intro", []))
    if "长城" in blob:
        return "cover"
    html = d.get("city_html", "")
    m = re.search(r"c-(\w+)'>\s*([^<]+?)\s*</span>", html)
    if m:
        city = m.group(2).strip()
        if city in photo_map: return photo_map[city]
    for c, ph in photo_map.items():
        if c in html: return ph
    return "cover"

def day_title(d):
    html = d.get("city_html", "")
    m = re.search(r"c-(\w+)'>\s*([^<]+?)\s*</span>\s*([^<]*)", html)
    if m:
        return re.sub(r"[<>/]", " ", m.group(2) + (m.group(3) or "")).strip() or d.get("stay","")
    return d.get("stay", d["day"])

def html_to_text(html):
    return re.sub(r"[<>/]", " ", html or "").strip()

def build_route(days, photo_map):
    """按城市段合并出路线总览：城市  起始–结束  标签。"""
    segs = []
    for d in days:
        city = None
        m = re.search(r"c-(\w+)'>\s*([^<]+?)\s*</span>", d.get("city_html",""))
        if m: city = m.group(2).strip()
        else: city = d.get("stay", d["day"]).split("·")[0].strip() or d["day"]
        date = d["date"]
        tag = ""
        if d.get("intro"):
            tag = d["intro"][0][0]
        if segs and segs[-1]["city"] == city:
            segs[-1]["end"] = date
        else:
            segs.append({"city": city, "start": date, "end": date, "tag": tag})
    lines = []
    for s in segs:
        rng = s["start"] if s["start"] == s["end"] else f"{s['start']}–{s['end']}"
        lines.append(f"{s['city']:<6} {rng}   {s['tag']}")
    return lines

# ---------------- 主流程 ----------------
def main():
    global prs, SW, SH, BLANK, PHOTO
    ap = argparse.ArgumentParser()
    ap.add_argument("data"); ap.add_argument("out")
    ap.add_argument("--photos", default=None)
    ap.add_argument("--budget", default=None)
    ap.add_argument("--cover-slug", default="cover")
    a = ap.parse_args()

    mod = load_data(a.data)
    datadir = os.path.dirname(os.path.abspath(a.data))
    PHOTO = a.photos or os.path.join(datadir, "ppt_assets", "photos")
    os.makedirs(PHOTO, exist_ok=True)

    days   = getattr(mod, "days", [])
    hotels = getattr(mod, "hotels", [])
    timeline= getattr(mod, "timeline", [])
    budget_data = getattr(mod, "budget", None)
    photo_map = dict(CITY_PHOTO); photo_map.update(getattr(mod, "photo_map", {}))
    TRIP_TITLE = getattr(mod, "TRIP_TITLE", "旅行行程")
    TRIP_SUB   = getattr(mod, "TRIP_SUB", "")
    CLOSING    = getattr(mod, "CLOSING", "行程圆满")

    # 预算（--budget xlsx 优先，其次数据模块 budget）
    budget = []
    if a.budget and os.path.exists(a.budget):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(a.budget)
            ws = wb["总览"]
            SKIP = {"总计","人均","建议区间(2人)","单人估算","2人总预算","备注","说明"}
            grand = 0
            for rr in range(5, 20):
                v = ws.cell(rr,1).value
                if v is None: continue
                if v == "总计": grand = _num(ws.cell(rr,2).value); continue
                pct = ws.cell(rr,3).value; b2 = ws.cell(rr,2).value
                if (v not in SKIP and b2 is not None and pct not in (None,"") and "%" in str(pct)
                        and not isinstance(v,(int,float))):
                    budget.append((v, _num(b2), pct))
        except Exception as e:
            print("预算读取失败:", e)
    elif budget_data:
        budget = [(b[0], int(b[1]), b[2]) for b in budget_data]

    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    BLANK = prs.slide_layouts[6]

    # 1. 封面
    s = new_slide(photo(a.cover_slug), 0.46)
    header_line(s, Inches(0.85), Inches(0.62), "TRAVEL ITINERARY")
    text(s, Inches(0.85), Inches(1.15), Inches(11.6), Inches(1.2), TRIP_TITLE, size=40, bold=True)
    if TRIP_SUB:
        text(s, Inches(0.87), Inches(2.35), Inches(11.6), Inches(0.6), TRIP_SUB, size=18, color=GRAY)
    sub2 = "%s  ·  %d 天 · %d 城" % (days[0]["date"] if days else "", len(days), len(set(build_route(days, photo_map))))
    text(s, Inches(0.85), Inches(6.9), Inches(11.6), Inches(0.5), sub2, size=14, color=GRAY)

    # 2. 路线总览
    s = new_slide(photo("transport") if os.path.exists(os.path.join(PHOTO,"transport.jpg")) else photo(a.cover_slug), 0.55)
    header_line(s, Inches(0.85), Inches(0.7), "ROUTE OVERVIEW")
    text(s, Inches(0.85), Inches(1.05), Inches(11), Inches(0.9), "由北向南 · 无折返", size=34, bold=True)
    route = build_route(days, photo_map)
    li, ri = route[: (len(route)+1)//2], route[(len(route)+1)//2:]
    col_bullets(s, Inches(0.9), Inches(2.1), Inches(6.0), Inches(4.8), li, "行程", size=18, gap=13)
    col_bullets(s, Inches(7.1), Inches(2.1), Inches(5.5), Inches(4.8), ri, "", size=18, gap=13)

    # 3. 关键数据（天数/城市数/预算）
    s = new_slide(photo("changbai") if os.path.exists(os.path.join(PHOTO,"changbai.jpg")) else photo(a.cover_slug), 0.55)
    header_line(s, Inches(0.85), Inches(0.7), "KEY FACTS")
    text(s, Inches(0.85), Inches(1.05), Inches(11), Inches(0.9), "行程关键数据", size=34, bold=True)
    facts = ["总天数：%d 天" % len(days),
             "途经城市：%d 座" % len(set(build_route(days, photo_map))),
             "每日单独成页：是",
             "背景图：彩色实景（非黑白）"]
    if budget:
        tot = sum(b[1] for b in budget)
        facts.append("总预算：¥%s（%d 类目）" % (f"{tot:,}", len(budget)))
    col_bullets(s, Inches(0.9), Inches(2.1), Inches(11.6), Inches(4.8), facts, "", size=20, gap=14)

    # 4. 每日页
    for d in days:
        ph = day_photo(d, photo_map)
        s = new_slide(photo(ph), 0.52)
        dayno = d["day"].replace("Day ", "")
        header_line(s, Inches(0.85), Inches(0.62), "DAY " + dayno + "  ·  " + d["date"] + " " + d["wd"])
        text(s, Inches(0.85), Inches(0.98), Inches(11.6), Inches(0.9), day_title(d), size=34, bold=True)
        panel(s, Inches(0.6), Inches(2.1), Inches(12.13), Inches(5.0), opacity=0.46)
        it_items = [shorten(it, 54) for it in d["itinerary"]]
        col_bullets(s, Inches(0.9), Inches(2.28), Inches(7.0), Inches(4.7), it_items, "当天行程", size=13.5, gap=7)
        # 右栏：必看景点 · 亮点
        attr_items = [shorten(f"{n}：{desc}", 34) for n, desc in d["intro"] if n != "—"]
        col_bullets(s, Inches(8.0), Inches(2.28), Inches(4.8), Inches(4.0), attr_items, "必看景点 · 亮点", size=11.5, gap=4)
        warns = [c for c in d["caution"] if "⚠" in c]
        warns = (warns + [c for c in d["caution"] if "⚠" not in c])[:2]
        if warns:
            text(s, Inches(8.0), Inches(6.4), Inches(4.8), Inches(0.42),
                 "⚠ " + " ｜ ".join(shorten(w, 30) for w in warns), size=10.5, color=WARN, spacing=1.0)
        stay_txt = d.get("stay", "")
        if stay_txt and "无住宿" not in stay_txt:
            text(s, Inches(0.9), Inches(6.95), Inches(11.6), Inches(0.5), "住宿：" + shorten(stay_txt, 120), size=13, color=GRAY)

    # 5. 预算构成
    if budget:
        s = new_slide(photo("shenyang") if os.path.exists(os.path.join(PHOTO,"shenyang.jpg")) else photo(a.cover_slug), 0.6)
        header_line(s, Inches(0.85), Inches(0.7), "BUDGET BREAKDOWN")
        tot = sum(b[1] for b in budget)
        text(s, Inches(0.85), Inches(1.05), Inches(11), Inches(0.9), "总预算 ¥%s" % f"{tot:,}", size=34, bold=True)
        maxpct = max(float(str(b[2]).replace("%","")) for b in budget if b[2]) or 1
        top = 2.5; rowH = 0.82; barL = Inches(3.4); barMax = Inches(6.4); barH = Inches(0.42)
        for i, (bname, amt, pct) in enumerate(budget):
            y = Inches(top + i*rowH)
            p = float(str(pct).replace("%","")) if pct else 0
            text(s, Inches(0.9), y, Inches(2.4), Inches(0.5), bname, size=17, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            tr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, barL, y+Inches(0.04), barMax, barH)
            tr.fill.solid(); tr.fill.fore_color.rgb = RGBColor(0x33,0x3D,0x47); tr.line.fill.background(); tr.shadow.inherit = False
            fw = int(barMax * (p/maxpct))
            fb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, barL, y+Inches(0.04), Emu(fw), barH)
            fb.fill.solid(); fb.fill.fore_color.rgb = BAR; fb.line.fill.background(); fb.shadow.inherit = False
            text(s, barL+Emu(fw)+Inches(0.12), y, Inches(2.4), Inches(0.5), "¥%s · %s" % (f"{int(amt):,}", pct), size=15, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)

    # 6. 住宿总览
    if hotels:
        s = new_slide(photo("dandong") if os.path.exists(os.path.join(PHOTO,"dandong.jpg")) else photo(a.cover_slug), 0.6)
        header_line(s, Inches(0.85), Inches(0.7), "STAY · 住宿总览")
        text(s, Inches(0.85), Inches(1.05), Inches(11.4), Inches(0.9), "住宿安排", size=30, bold=True)
        ty = Inches(2.2); rowH = Inches(0.5)
        for i, row in enumerate(hotels):
            city, nights, status, hname, room, reason, note = (list(row) + ["","","","","","",""])[:7]
            y = ty + Emu(int(rowH)*i)
            if i % 2 == 0:
                bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y, Inches(11.9), rowH)
                bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0,0,0); bg.line.fill.background()
                srgb = bg.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
                srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(0.4*65535))}))
                bg.shadow.inherit = False
            st_color = ST_OK if "✅" in status else ST_REC
            text(s, Inches(0.85), y, Inches(1.7), rowH, status, size=12.5, bold=True, color=st_color, anchor=MSO_ANCHOR.MIDDLE)
            text(s, Inches(2.5), y, Inches(1.9), rowH, city, size=12.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
            text(s, Inches(4.4), y, Inches(5.4), rowH, shorten(hname, 46), size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
            text(s, Inches(9.9), y, Inches(2.7), rowH, shorten(note, 34), size=11.5, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)

    # 7. 预约时间轴
    if timeline:
        s = new_slide(photo("beijing") if os.path.exists(os.path.join(PHOTO,"beijing.jpg")) else photo(a.cover_slug), 0.58)
        header_line(s, Inches(0.85), Inches(0.7), "BEFORE YOU GO · 预约时间轴")
        text(s, Inches(0.85), Inches(1.05), Inches(11.4), Inches(0.9), "行前务必提前实名预约", size=30, bold=True)
        tl = [f"{p} {item} — {when}：{desc}" for p, item, when, desc in timeline]
        col_bullets(s, Inches(0.9), Inches(2.1), Inches(11.6), Inches(4.9), tl, "", size=14.5, gap=9)

    # 8. 收尾
    s = new_slide(photo("taishan") if os.path.exists(os.path.join(PHOTO,"taishan.jpg")) else photo(a.cover_slug), 0.5)
    text(s, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.4), CLOSING, size=44, bold=True, align=PP_ALIGN.CENTER)
    text(s, Inches(1.0), Inches(4.4), Inches(11.3), Inches(0.7),
         "%s 天 · %d 城 · 一路向北再到南" % (len(days), len(set(build_route(days, photo_map)))),
         size=18, color=GRAY, align=PP_ALIGN.CENTER)

    prs.save(a.out)
    print("PPTX saved:", a.out, "| slides:", len(prs.slides._sldIdLst))

def _num(v):
    if v is None: return 0
    if isinstance(v, (int, float)): return int(v)
    d = re.sub(r"[^\d]", "", str(v)); return int(d) if d else 0

if __name__ == "__main__":
    main()
