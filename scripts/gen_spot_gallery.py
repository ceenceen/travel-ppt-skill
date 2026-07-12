# -*- coding: utf-8 -*-
"""模式 C · 景点图鉴：把 GROUPS + DAYS 渲染成深色高级网格 PPTX（travel-ppt 视觉体系）。

用法 / Usage:
  python gen_spot_gallery.py --data spot_gallery_data.py --out 景点图鉴.pptx

数据模块需定义 / The data module must define:
  GROUPS : [(cn_theme, en_sub, [spot_cn, ...]), ...]
  DAYS   : {spot_cn: int}   # 可选，用于每个景点名旁的「Day N」标签

版式：深色背景 + 白色强调（原橘色图标已改白）+ 微软雅黑；杂志式紧凑网格，
图片按单元格比例 cover 裁切填满（无白边），底部半透明黑条 + 白色小方块 + 白字可编辑名称 + Day N 标签。
图片来自 {IMG_DIR}/{景点名}.jpg（由 fetch_spot_photos.py 抓取，优先 Pexels、缺图 AI 兜底）。
"""
import os
import importlib.util
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

IMG_DIR = 'spot_gallery'
CROP_DIR = 'spot_gallery/_cropped'
os.makedirs(CROP_DIR, exist_ok=True)

# 配色（travel-ppt 视觉体系；白色强调，非橘）
BK = RGBColor(0x0B, 0x0E, 0x14)      # 深色背景
PANEL = RGBColor(0x15, 0x1A, 0x22)   # 占位面板
W = RGBColor(0xFF, 0xFF, 0xFF)
SUB = RGBColor(0xB8, 0xC0, 0xCC)
ACCENT = RGBColor(0xFF, 0xFF, 0xFF)  # 白色强调（原橘色 E86D4E）
LINE = RGBColor(0x2A, 0x31, 0x3E)


def load_data(path):
    spec = importlib.util.spec_from_file_location('spot_data', path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BK


def set_alpha(shape, pct):
    sp = shape.fill._xPr.find(qn('a:solidFill'))
    if sp is None:
        return
    clr = sp.find(qn('a:srgbClr'))
    if clr is None:
        return
    a = clr.find(qn('a:alpha'))
    if a is None:
        a = clr.makeelement(qn('a:alpha'), {})
        clr.append(a)
    a.set('val', str(int(pct * 1000)))


def cover_crop(src, dst, target_ratio):
    im = Image.open(src).convert('RGB')
    iw, ih = im.size
    r = iw / ih
    if r > target_ratio:
        nw = int(ih * target_ratio)
        x0 = (iw - nw) // 2
        im = im.crop((x0, 0, x0 + nw, ih))
    else:
        nh = int(iw / target_ratio)
        y0 = (ih - nh) // 2
        im = im.crop((0, y0, iw, y0 + nh))
    im.save(dst, quality=88)


def build_grid(cn_name, en_name, spots, days):
    prs = build_grid.prs
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)

    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.30), Inches(9.5), Inches(0.62))
    tf = tb.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = cn_name
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = W; r.font.name = 'Microsoft YaHei'

    en = s.shapes.add_textbox(Inches(0.58), Inches(0.86), Inches(9.5), Inches(0.32))
    er = en.text_frame.paragraphs[0].add_run()
    er.text = en_name
    er.font.size = Pt(11); er.font.color.rgb = SUB; er.font.name = 'Microsoft YaHei'

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(0.80), Inches(0.9), Inches(0.045))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False

    cnt = s.shapes.add_textbox(Inches(11.3), Inches(0.38), Inches(1.7), Inches(0.5))
    cr = cnt.text_frame.paragraphs[0]
    cr.alignment = PP_ALIGN.RIGHT
    crn = cr.add_run()
    crn.text = f'{len(spots)} 处'
    crn.font.size = Pt(13); crn.font.bold = True; crn.font.color.rgb = ACCENT; crn.font.name = 'Microsoft YaHei'

    n = len(spots)
    cols = 4 if n >= 8 else 3
    rows = (n + cols - 1) // cols
    L = 0.55; Rm = 0.55; T = 1.30; B = 0.30; gut = 0.12
    grid_w = 13.333 - L - Rm
    grid_h = 7.5 - T - B
    cell_w = (grid_w - gut * (cols - 1)) / cols
    cell_h = (grid_h - gut * (rows - 1)) / rows
    target_ratio = cell_w / cell_h

    for i, cn in enumerate(spots):
        c = i % cols
        rr = i // cols
        cx = L + c * (cell_w + gut)
        cy = T + rr * (cell_h + gut)

        src = os.path.join(IMG_DIR, cn + '.jpg')
        if os.path.exists(src):
            dst = os.path.join(CROP_DIR, f'{cn}.jpg')
            try:
                cover_crop(src, dst, target_ratio)
                s.shapes.add_picture(dst, Inches(cx), Inches(cy), Inches(cell_w), Inches(cell_h))
            except Exception:
                pass
        else:
            ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(cy), Inches(cell_w), Inches(cell_h))
            ph.fill.solid(); ph.fill.fore_color.rgb = PANEL; ph.line.fill.background(); ph.shadow.inherit = False

        bd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(cy), Inches(cell_w), Inches(cell_h))
        bd.fill.background(); bd.line.color.rgb = LINE; bd.line.width = Pt(0.75); bd.shadow.inherit = False

        strip_h = 0.34
        sy = cy + cell_h - strip_h
        strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(sy), Inches(cell_w), Inches(strip_h))
        strip.fill.solid(); strip.fill.fore_color.rgb = RGBColor(0x0A, 0x0C, 0x10)
        set_alpha(strip, 58); strip.line.fill.background(); strip.shadow.inherit = False

        dot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx + 0.10), Inches(sy + strip_h / 2 - 0.055), Inches(0.06), Inches(0.11))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT; dot.line.fill.background(); dot.shadow.inherit = False

        nt = s.shapes.add_textbox(Inches(cx + 0.22), Inches(sy), Inches(cell_w - 0.28), Inches(strip_h))
        ntf = nt.text_frame
        ntf.word_wrap = True
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ntf.margin_left = 0; ntf.margin_right = 0; ntf.margin_top = 0; ntf.margin_bottom = 0
        np = ntf.paragraphs[0]
        np.alignment = PP_ALIGN.LEFT
        day = days.get(cn)
        if day is not None:
            dr = np.add_run()
            dr.text = f'Day {day}  '
            dr.font.size = Pt(9); dr.font.bold = False; dr.font.color.rgb = SUB; dr.font.name = 'Microsoft YaHei'
        nr = np.add_run()
        nr.text = cn
        nr.font.size = Pt(12 if len(cn) <= 6 else 10.5)
        nr.font.bold = True; nr.font.color.rgb = W; nr.font.name = 'Microsoft YaHei'


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--data', required=True)
    ap.add_argument('--out', default='景点图鉴.pptx')
    a = ap.parse_args()

    mod = load_data(a.data)
    days = getattr(mod, 'DAYS', {})
    GROUPS = getattr(mod, 'GROUPS', [])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build_grid.prs = prs

    for cn_name, en_name, spots in GROUPS:
        build_grid(cn_name, en_name, spots, days)
    prs.save(a.out)
    print(f'SAVED {a.out}  slides={len(prs.slides._sldIdLst)}')


if __name__ == '__main__':
    main()
