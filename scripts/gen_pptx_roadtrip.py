"""gen_pptx_roadtrip.py —— 自驾/户外环线 每日行程 PPTX 生成器（双框排版）

每页版式（16:9, 13.33x7.5in）：
  · 全幅风景照背景（cover 铺满 + 暗化 overlay，保证文字可读）
  · 右上角「上下双框」：上 = 真实路线地图（gen_maps.py 产出），下 = 当天景点图
      - 景点图按原比例 contain 居中（不拉伸），留白用半透明深底盖住
      - 景点图右上角带当天主要景点名标注
  · 左列四块文字：行程说明 / 经典介绍 / 注意事项 / 住宿点

依赖: python-pptx, Pillow
数据来源: trip_data.py（与本 skill 的 gen_maps.py 共用），需提供：
    S  —— 页面序列，元素为 (类型, 数据dict)
         类型见 builders：cover / section / stats / day / tips / closing
    get_img(idx)        —— 返回通用背景风光照路径（用于封面/分隔/数据页）
    get_spot(day)       —— 返回当天景点图路径（AI 生成写实图，见 SKILL.md 图片管线）
    get_map(day)        —— 返回当天地图路径（maps/day_map_{day}.png）
  每日 day 页数据 dict 示例：
    {'img': get_spot(1), 'bg': get_img(6), 'map': get_map(1),
     'day_label': 'DAY 1 · 8.10', 'route': '拉萨 → 日喀则',
     'km': '360', 'time': '6h', 'highest_alt': '3836m',
     'trip': [...], 'intro': [...], 'caution': [...], 'stay': [...]}

用法: 在含 trip_data.py 的旅行目录下执行
    python gen_pptx_roadtrip.py
输出: <DECK_DIR>/<TRIP_TITLE>-roadtrip.pptx（被占用时自动换名）
"""
import sys, os

sys.path.insert(0, r'C:/Users/Administrator/AppData/Roaming/Python/Python313/site-packages')
sys.path.insert(0, os.getcwd())
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from trip_data import S, get_img, get_spot, get_map, TRIP_TITLE, DECK_DIR

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

W = RGBColor(255, 255, 255)
G = RGBColor(167, 167, 167)
DG = RGBColor(111, 111, 111)
VL = RGBColor(242, 242, 242)
BK = RGBColor(0, 0, 0)
ACCENT = RGBColor(232, 109, 78)

def add_bg(slide, img_path, intensity=0.35):
    """全幅背景：cover 逻辑按比例铺满整页，超出部分居中裁切，再叠暗化层。"""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BK
    if img_path and os.path.exists(img_path):
        with Image.open(img_path) as im:
            iw, ih = im.size
        slide_w, slide_h = 13.33, 7.5
        if (iw / ih) >= (slide_w / slide_h):
            rh = slide_h
            rw = slide_h * iw / ih
            left = (slide_w - rw) / 2
            top = 0.0
        else:
            rw = slide_w
            rh = slide_w * ih / iw
            top = (slide_h - rh) / 2
            left = 0.0
        slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                                 Inches(rw), Inches(rh))
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BK
    shape.line.fill.background()
    sp_pr = shape._element.find(qn('p:spPr'))
    if sp_pr is None:
        sp_pr = shape._element.makeelement(qn('p:spPr'), {})
        shape._element.insert(0, sp_pr)
    for sf in sp_pr.findall(qn('a:solidFill')):
        sp_pr.remove(sf)
    solid_fill = sp_pr.makeelement(qn('a:solidFill'), {})
    srgb_clr = solid_fill.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alpha = srgb_clr.makeelement(qn('a:alpha'), {'val': str(int((1-intensity)*65535))})
    srgb_clr.append(alpha)
    solid_fill.append(srgb_clr)
    sp_pr.append(solid_fill)

def tb(slide, left, top, width, height, texts, font_name='Microsoft YaHei', alignment=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, clr, bld) in enumerate(texts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment or PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = clr
        run.font.bold = bld
        run.font.name = font_name
        if i > 0:
            p.space_before = Pt(10)
    return txBox

def block(slide, x, y, w, title, lines, fs=15.5, step=0.34):
    """小节：标题（强调色）+ 多行内容（浅色），可调字号与行高。"""
    tb(slide, Inches(x), Inches(y), Inches(w), Inches(0.4), [(title, 16, ACCENT, True)])
    yy = y + 0.42
    for ln in lines:
        tb(slide, Inches(x), Inches(yy), Inches(w), Inches(step + 0.16), [(ln, fs, VL, False)])
        yy += step

def build_cover(data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, data['img'], 0.35)
    tb(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.5),
       [(data['word'], 84, W, True), (data['sub'], 28, W, False)],
       alignment=PP_ALIGN.CENTER)
    tb(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
       [(data['meta'], 18, DG, False)], alignment=PP_ALIGN.CENTER)

def build_section(data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, data['img'], 0.5)
    tb(slide, Inches(0.8), Inches(0.6), Inches(6), Inches(0.4),
       [(data['label'], 14, DG, False)])
    tb(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1.5),
       [(data['title'], 64, W, True)])
    if 'body' in data:
        tb(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(1),
           [(data['body'], 22, G, False)])

def build_stats(data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, data['img'], 0.45)
    tb(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.8),
       [(data['hl'], 42, W, True)])
    y = 1.6
    for item in data['stats']:
        if len(item) == 3:
            num, unit, label = item
        else:
            num, label = item[0], item[1]
            unit = ''
        tb(slide, Inches(0.6), Inches(y), Inches(5), Inches(0.8),
           [(f'{num}{unit}', 72, W, True)])
        tb(slide, Inches(0.6), Inches(y+1.0), Inches(5), Inches(0.5),
           [(label, 18, G, False)])
        y += 1.7

def _set_alpha(shape, alpha_frac):
    """给黑色填充形状设置透明度 (0=全透,1=不透)。"""
    sp_pr = shape._element.find(qn('p:spPr'))
    if sp_pr is None:
        sp_pr = shape._element.makeelement(qn('p:spPr'), {})
        shape._element.insert(0, sp_pr)
    for sf in sp_pr.findall(qn('a:solidFill')):
        sp_pr.remove(sf)
    solid = sp_pr.makeelement(qn('a:solidFill'), {})
    srgb = solid.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alpha = srgb.makeelement(qn('a:alpha'), {'val': str(int((1 - alpha_frac) * 65535))})
    srgb.append(alpha); solid.append(srgb); sp_pr.append(solid)

def build_day(data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 全幅风景照背景（cover 铺满 + 暗化，保证文字/小框可读）
    add_bg(slide, data.get('bg') or data.get('img'), 0.62)
    # 顶部标题区
    tb(slide, Inches(0.8), Inches(0.45), Inches(6), Inches(0.4),
       [(data['day_label'], 14, DG, False)])
    tb(slide, Inches(0.8), Inches(0.82), Inches(11.8), Inches(1.0),
       [(data['route'], 38, W, True)])
    # 数据条：里程 / 时长 / 最高海拔
    strip = [('里程 DISTANCE', data['km'] + ' km'),
             ('时长 TIME', data['time']),
             ('最高海拔 ALT', data['highest_alt'])]
    sx = 0.8
    for lbl, val in strip:
        tb(slide, Inches(sx), Inches(1.95), Inches(2.2), Inches(0.35),
           [(lbl, 11, DG, False)])
        tb(slide, Inches(sx), Inches(2.28), Inches(2.2), Inches(0.5),
           [(val, 24, W, True)])
        sx += 2.2
    # ── 左列：行程说明 / 经典介绍 / 注意事项 / 住宿点（4 块紧凑）
    lw = 5.8
    block(slide, 0.8, 3.05, lw, '▸ 行程说明', data['trip'], fs=12.5, step=0.21)
    block(slide, 0.8, 4.30, lw, '▸ 经典介绍', data['intro'], fs=12.5, step=0.21)
    block(slide, 0.8, 5.50, lw, '▸ 注意事项', data['caution'], fs=12.5, step=0.21)
    block(slide, 0.8, 6.55, lw, '▸ 住宿点', data['stay'], fs=12.5, step=0.21)
    # ── 右上角上下双框：上=真实路线地图，下=今日景点图
    # 上框：真实路线地图
    mp = data.get('map')
    if mp and os.path.exists(mp):
        slide.shapes.add_picture(mp, Inches(6.9), Inches(3.05), Inches(5.9), Inches(2.0))
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(6.9), Inches(3.05), Inches(5.9), Inches(2.0))
        border.fill.background()
        border.line.color.rgb = RGBColor(100, 100, 100)
        border.line.width = Pt(0.75)
    # 下框：今日景点图（保持原比例，contain 居中，不拉伸）
    ip = data.get('img')
    if ip and os.path.exists(ip):
        box_x, box_y, box_w, box_h = 6.9, 5.15, 5.9, 2.0
        with Image.open(ip) as _im:
            iw, ih = _im.size
        # contain：按原比例完整显示，留白区域用半透明深底盖住
        scale = min(box_w / iw, box_h / ih)
        dw, dh = iw * scale, ih * scale
        dx = box_x + (box_w - dw) / 2
        dy = box_y + (box_h - dh) / 2
        back = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(box_x), Inches(box_y),
                                      Inches(box_w), Inches(box_h))
        back.fill.solid()
        back.fill.fore_color.rgb = BK
        back.line.fill.background()
        _set_alpha(back, 0.55)
        slide.shapes.add_picture(ip, Inches(dx), Inches(dy), Inches(dw), Inches(dh))
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(dx), Inches(dy), Inches(dw), Inches(dh))
        border.fill.background()
        border.line.color.rgb = RGBColor(100, 100, 100)
        border.line.width = Pt(0.75)
        # 景点图右上角标注：当天主要景点名
        spot = data['intro'][0].split('：')[0].split('（')[0].split('(')[0]
        lab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(dx + dw - 2.7), Inches(dy + 0.10),
                                     Inches(2.7), Inches(0.42))
        lab.fill.solid()
        lab.fill.fore_color.rgb = BK
        lab.line.fill.background()
        _set_alpha(lab, 0.6)
        tf = lab.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = spot
        run.font.size = Pt(12)
        run.font.color.rgb = W
        run.font.bold = True
        run.font.name = 'Microsoft YaHei'

def build_tips(data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, data.get('img', ''), 0.5)
    tb(slide, Inches(0.8), Inches(0.6), Inches(6), Inches(0.8),
       [(data['hl'], 36, W, True)])
    y = 1.8
    for item in data['items']:
        tb(slide, Inches(0.8), Inches(y), Inches(11), Inches(0.5),
           [(item, 20, VL, False)])
        y += 0.55

def build_closing(data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, data['img'], 0.7)
    tb(slide, Inches(2), Inches(3.0), Inches(9), Inches(1.5),
       [(data['title'], 72, W, True)], alignment=PP_ALIGN.CENTER)
    tb(slide, Inches(2), Inches(4.5), Inches(9), Inches(0.8),
       [(data['sub'], 24, DG, False)], alignment=PP_ALIGN.CENTER)

builders = {
    'cover': build_cover,
    'section': build_section,
    'stats': build_stats,
    'day': build_day,
    'tips': build_tips,
    'closing': build_closing
}

for stype, data in S:
    builders[stype](data)

output_path = os.path.join(DECK_DIR, f'{TRIP_TITLE}-roadtrip.pptx')
try:
    prs.save(output_path)
except PermissionError:
    output_path = os.path.join(DECK_DIR, f'{TRIP_TITLE}-roadtrip-alt.pptx')
    prs.save(output_path)
fsize = os.path.getsize(output_path) / 1024 / 1024
print(f'PPTX saved: {output_path}')
print(f'File size: {fsize:.1f} MB')
print(f'Total slides: {len(S)}')
